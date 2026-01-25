import os
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

def get_shape_text_and_style(shape):
    """Extracts text and basic style info from a shape."""
    if not shape.has_text_frame:
        return None
    
    data = []
    for paragraph in shape.text_frame.paragraphs:
        p_text = paragraph.text.strip()
        if not p_text:
            continue
        
        # Get style from the first run if available, else paragraph, else defaults
        font_name = "Default"
        font_size = "Default"
        is_bold = False
        is_italic = False
        
        if paragraph.runs:
            run = paragraph.runs[0]
            if run.font.name:
                font_name = run.font.name
            if run.font.size:
                font_size = run.font.size.pt
            if run.font.bold:
                is_bold = True
            if run.font.italic:
                is_italic = True
        
        data.append({
            'text': p_text,
            'font_name': font_name,
            'font_size': font_size,
            'bold': is_bold,
            'italic': is_italic,
            'left': int(shape.left.pt) if hasattr(shape, 'left') else 0,
            'top': int(shape.top.pt) if hasattr(shape, 'top') else 0,
            'width': int(shape.width.pt) if hasattr(shape, 'width') else 0,
            'height': int(shape.height.pt) if hasattr(shape, 'height') else 0
        })
    return data

def analyze_ppt(internal_path, benchmark_path):
    print("Loading PPT files...")
    try:
        prs_internal = Presentation(internal_path)
        prs_benchmark = Presentation(benchmark_path)
    except Exception as e:
        print(f"Error loading files: {e}")
        return

    layout_comparison = []
    translation_comparison = []
    
    # Iterate through slides (assuming same count/order for direct comparison)
    max_slides = min(len(prs_internal.slides), len(prs_benchmark.slides))
    
    for i in range(max_slides):
        slide_int = prs_internal.slides[i]
        slide_ben = prs_benchmark.slides[i]
        
        shapes_int = [s for s in slide_int.shapes if s.has_text_frame]
        shapes_ben = [s for s in slide_ben.shapes if s.has_text_frame]
        
        # Sort shapes by position (top-left) to try and match them up
        shapes_int.sort(key=lambda s: (s.top, s.left))
        shapes_ben.sort(key=lambda s: (s.top, s.left))
        
        # Simple matching by index for now
        for j in range(min(len(shapes_int), len(shapes_ben))):
            s_int = shapes_int[j]
            s_ben = shapes_ben[j]
            
            data_int_list = get_shape_text_and_style(s_int)
            data_ben_list = get_shape_text_and_style(s_ben)
            
            if not data_int_list or not data_ben_list:
                continue

            # Compare first paragraph for style checking to avoid noise
            d_int = data_int_list[0]
            d_ben = data_ben_list[0]
            
            # Layout Comparison
            diffs = []
            if d_int['font_name'] != d_ben['font_name']:
                diffs.append(f"Font: {d_int['font_name']} vs {d_ben['font_name']}")
            if d_int['font_size'] != d_ben['font_size']:
                 # Tolerate small differences
                 if isinstance(d_int['font_size'], (int, float)) and isinstance(d_ben['font_size'], (int, float)):
                     if abs(d_int['font_size'] - d_ben['font_size']) > 1:
                         diffs.append(f"Size: {d_int['font_size']} vs {d_ben['font_size']}")
            if d_int['bold'] != d_ben['bold']:
                diffs.append(f"Bold: {d_int['bold']} vs {d_ben['bold']}")
            
            # Position diff
            pos_diff = abs(d_int['left'] - d_ben['left']) + abs(d_int['top'] - d_ben['top'])
            if pos_diff > 20: # Tolerance
                diffs.append(f"Pos Diff: >20pt")

            if diffs:
                layout_comparison.append({
                    'slide': i + 1,
                    'item': f"Text Block {j+1}",
                    'internal': f"Font: {d_int['font_name']}, Size: {d_int['font_size']}",
                    'benchmark': f"Font: {d_ben['font_name']}, Size: {d_ben['font_size']}",
                    'diff': ", ".join(diffs)
                })

            # Translation Comparison
            # Join all text in shape
            text_int = "\n".join([d['text'] for d in data_int_list])
            text_ben = "\n".join([d['text'] for d in data_ben_list])
            
            # Simple check if text is significantly different
            if text_int != text_ben:
                translation_comparison.append({
                    'slide': i + 1,
                    'original': "(See original)", # Placeholder as we are comparing Output vs Benchmark
                    'internal': text_int[:50].replace('\n', ' ') + "..." if len(text_int) > 50 else text_int.replace('\n', ' '),
                    'benchmark': text_ben[:50].replace('\n', ' ') + "..." if len(text_ben) > 50 else text_ben.replace('\n', ' '),
                    'issue': "Content mismatch"
                })

    # Generate Markdown Report
    output_file = "analysis_report.md"
    with open(output_file, "w", encoding="utf-8") as f:
        f.write("## 분석 결과\n")
        f.write("\n### 1. 레이아웃 비교\n")
        f.write("| 슬라이드 | 항목 | RCA_내버전 | RCA_유료버전 | 차이점 |\n")
        f.write("|---|---|---|---|---|\n")
        for row in layout_comparison[:20]: # Limit rows
            f.write(f"| {row['slide']} | {row['item']} | {row['internal']} | {row['benchmark']} | {row['diff']} |\n")
        if not layout_comparison:
            f.write("| - | - | - | - | 차이점 없음 |\n")

        f.write("\n### 2. 번역 품질 비교\n")
        f.write("| 슬라이드 | 내버전 번역 | 유료버전 번역 | 문제점 |\n")
        f.write("|---|---|---|---|\n")
        for row in translation_comparison[:20]:
             f.write(f"| {row['slide']} | {row['internal']} | {row['benchmark']} | {row['issue']} |\n")
        if not translation_comparison:
             f.write("| - | - | - | 차이점 없음 |\n")

        f.write("\n### 3. 개선 우선순위 To-Do List\n")
        f.write("- [ ] (이곳은 분석 결과를 바탕으로 채워주세요)\n")
    
    print(f"Analysis saved to {output_file}")

if __name__ == "__main__":
    base_dir = r"c:\Users\nomus\Desktop\구글 동기화\안티그래비티\테스트\project\2. ppt"
    internal = os.path.join(base_dir, "2. RCA 보고서_내버전(v3).pptx")
    benchmark = os.path.join(base_dir, "3. RCA 보고서_유료버전.pptx")
    
    if not os.path.exists(internal):
        print(f"File not found: {internal}")
    elif not os.path.exists(benchmark):
        print(f"File not found: {benchmark}")
    else:
        analyze_ppt(internal, benchmark)
