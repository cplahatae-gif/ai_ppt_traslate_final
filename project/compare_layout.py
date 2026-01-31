import sys
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def get_color_hex(color):
    """Extract color as hex string"""
    try:
        if color is None:
            return "None"
        if color.type == 1:  # RGB
            return str(color.rgb)
        elif color.type == 2:  # Theme
            return f"Theme:{color.theme_color}"
        else:
            return f"Type:{color.type}"
    except:
        return "Unknown"

def get_highlight(run):
    """Try to get highlight/background color if available"""
    try:
        # python-pptx doesn't directly support highlight, check for solidFill in rPr
        return "N/A"
    except:
        return "N/A"

def analyze_ppt_comprehensive(ppt_path, output_lines):
    prs = Presentation(ppt_path)
    output_lines.append(f"\n{'='*80}")
    output_lines.append(f"FILE: {ppt_path}")
    output_lines.append(f"{'='*80}\n")
    
    for slide_idx, slide in enumerate(prs.slides):
        output_lines.append(f"--- SLIDE {slide_idx + 1} ---\n")
        
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            text_frame = shape.text_frame
            shape_text = text_frame.text.strip()
            if not shape_text:
                continue
                
            output_lines.append(f"[Shape ID: {shape.shape_id}]")
            
            for p_idx, paragraph in enumerate(text_frame.paragraphs):
                para_text = "".join(run.text for run in paragraph.runs)
                if not para_text.strip():
                    continue
                    
                output_lines.append(f"  Paragraph {p_idx}: \"{para_text[:50]}{'...' if len(para_text)>50 else ''}\"")
                
                for r_idx, run in enumerate(paragraph.runs):
                    if not run.text.strip() and len(run.text) < 3:
                        continue  # Skip whitespace-only runs
                    
                    font = run.font
                    details = []
                    
                    # Font Name
                    if font.name:
                        details.append(f"Font:{font.name}")
                    
                    # Font Size
                    if font.size:
                        details.append(f"Size:{font.size.pt:.1f}pt")
                    
                    # Bold
                    details.append(f"Bold:{font.bold}")
                    
                    # Italic
                    details.append(f"Italic:{font.italic}")
                    
                    # Underline
                    details.append(f"Underline:{font.underline}")
                    
                    # Color
                    color_hex = get_color_hex(font.color)
                    details.append(f"Color:{color_hex}")
                    
                    # Strikethrough
                    if hasattr(font, 'strikethrough'):
                        details.append(f"Strike:{font.strikethrough}")
                    
                    output_lines.append(f"    Run {r_idx}: \"{run.text[:30]}{'...' if len(run.text)>30 else ''}\" | {' | '.join(details)}")
            
            output_lines.append("")  # Empty line between shapes

# Main execution
original_file = r"c:\Users\nomus\Desktop\구글 동기화\안티그래비티\테스트\project\ai ppt 번역기\2. ppt\test_file.pptx"
translated_file = r"c:\Users\nomus\Desktop\구글 동기화\안티그래비티\테스트\project\ai ppt 번역기\2. ppt\test_file_p1-2_en.pptx"
output_report = r"c:\Users\nomus\Desktop\구글 동기화\안티그래비티\테스트\project\ai ppt 번역기\2. ppt\layout_comparison_report.txt"

output_lines = []
output_lines.append("PPT LAYOUT COMPARISON REPORT")
output_lines.append("Generated for layout difference analysis\n")

try:
    analyze_ppt_comprehensive(original_file, output_lines)
    analyze_ppt_comprehensive(translated_file, output_lines)
    
    with open(output_report, "w", encoding="utf-8") as f:
        f.write("\n".join(output_lines))
    
    print(f"Report saved to: {output_report}")
except Exception as e:
    print(f"Error: {e}")
    import traceback
    traceback.print_exc()
