import sys
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt

def analyze_slide(ppt_path, slide_index, output_file):
    prs = Presentation(ppt_path)
    
    # Adjust slide_index (0-based)
    if slide_index >= len(prs.slides):
        print(f"Slide {slide_index + 1} not found in {ppt_path}")
        return

    slide = prs.slides[slide_index]
    
    with open(output_file, "a", encoding="utf-8") as f:
        f.write(f"\n--- Analyzing {ppt_path} Slide {slide_index + 1} ---\n")
        
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
                
            text_frame = shape.text_frame
            f.write(f"\n[Shape ID: {shape.shape_id}]\n")
            f.write(f"  Word Wrap: {text_frame.word_wrap}\n")
            
            for p_idx, paragraph in enumerate(text_frame.paragraphs):
                f.write(f"  Paragraph {p_idx}:\n")
                
                # Check for run level formatting
                runs_info = []
                for r_idx, run in enumerate(paragraph.runs):
                    text = run.text
                    font = run.font
                    color = str(font.color.rgb) if font.color and font.color.type == 1 else "Theme/None"
                    bold = font.bold
                    size = font.size.pt if font.size else "Inherit"
                    runs_info.append(f"Run {r_idx}: '{text}' (Bold: {bold}, Color: {color}, Size: {size})")
                
                for info in runs_info:
                    f.write(f"    {info}\n")

# Analyze Page 1 (Index 0) and Page 5 (Index 4)
original_file = "project/ai ppt 번역기/2. ppt/R2_Safety Leadership Guide_2602.pptx"
translated_file = "project/ai ppt 번역기/2. ppt/R2_Safety Leadership Guide_2602_p1-8_en.pptx"
output_report = "project/ai ppt 번역기/2. ppt/style_analysis_report.txt"

# Clear output file
open(output_report, "w").close()

try:
    analyze_slide(original_file, 0, output_report)
    analyze_slide(translated_file, 0, output_report)
    
    analyze_slide(original_file, 4, output_report)
    analyze_slide(translated_file, 4, output_report)
    
    print(f"Analysis saved to {output_report}")
except Exception as e:
    print(f"Error: {e}")
